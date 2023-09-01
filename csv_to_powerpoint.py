from pptx import Presentation
import csv


def main():
    # Input and output paths
    pptx_path = '../../powerpoint/data/output.pptx'
    csv_path = '../../powerpoint/data/output.csv'

    rows = []
    # Read in the data from the csv
    with open(csv_path, 'r', newline='', encoding='utf-8') as csvfile:
        csv_reader = csv.reader(csvfile)
        for row in csv_reader:
            rows.append(row)

    # Create a new presentation
    prs = Presentation()
    # Add a new slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    # Find the longest row
    longest_row = 0
    for row in rows:
        if len(row) > longest_row:
            longest_row = len(row)

    # Add a new table to the slide
    table = slide.shapes.add_table(rows=len(rows), cols=longest_row, left=0, top=0, width=prs.slide_width, height=prs.slide_height).table
    # Populate the table with the data from the csv
    for i in range(len(rows)):
        for j in range(len(rows[i])):
            table.cell(i, j).text = rows[i][j]
    # Save the presentation
    prs.save(pptx_path)


if __name__ == "__main__":
    main()
